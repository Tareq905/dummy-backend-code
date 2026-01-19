from pathlib import Path
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


def extract_text(file_path: Path) -> str:
    """
    Extract raw text from supported document formats.
    Returns empty string if extraction fails.
    """
    suffix = file_path.suffix.lower()

    try:
        if suffix == ".pdf":
            return _read_pdf(file_path)

        if suffix == ".docx":
            return _read_docx(file_path)

        if suffix == ".xlsx":
            return _read_xlsx(file_path)

        if suffix == ".pptx":
            return _read_pptx(file_path)

        if suffix in {".txt", ".csv"}:
            return _read_text_file(file_path)

    except Exception:
        # Never crash the pipeline for a single bad file
        return ""

    return ""


def _read_pdf(path: Path) -> str:
    text = []

    try:
        reader = PdfReader(path)
    except Exception:
        return ""

    for page in reader.pages:
        try:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
        except Exception:
            continue

    return "\n".join(text).strip()


def _read_docx(path: Path) -> str:
    try:
        doc = Document(path)
    except Exception:
        return ""

    return "\n".join(
        p.text.strip()
        for p in doc.paragraphs
        if p.text and p.text.strip()
    ).strip()


def _read_xlsx(path: Path) -> str:
    text = []

    try:
        wb = load_workbook(path, data_only=True)
    except Exception:
        return ""

    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    try:
                        text.append(str(cell))
                    except Exception:
                        continue

    return "\n".join(text).strip()


def _read_pptx(path: Path) -> str:
    text = []

    try:
        prs = Presentation(path)
    except Exception:
        return ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content = shape.text
                if content and content.strip():
                    text.append(content.strip())

    return "\n".join(text).strip()


def _read_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="ignore").strip()
    except Exception:
        return ""